"""
סוכן AI אחד (Agent) עם ממשק Streamlit, מבוסס Google Gemini (ה-API האישי שלך, לא תלוי במכללה).
כולל: קבצים, shell, Gmail, חיפוש אינטרנט, ויצירת קבצי Word/PowerPoint.
כשהמכסה היומית של Gemini מתרוקנת, הסוכן עובר אוטומטית לגיבוי - ה-API של המכללה (IAC).
"""

import base64
import io
import json
import os
import time
import uuid
from datetime import datetime
from email import encoders
from email.mime.base import MIMEBase
from email.mime.multipart import MIMEMultipart
from email.mime.text import MIMEText
from pathlib import Path

import requests
import streamlit as st
from dotenv import load_dotenv
from docx import Document
from pptx import Presentation
from pptx.util import Inches
from openpyxl import Workbook, load_workbook
from google import genai
from google.genai import errors as genai_errors
from google.genai import types

from google.auth.transport.requests import Request
from google.oauth2.credentials import Credentials
from google_auth_oauthlib.flow import InstalledAppFlow
from googleapiclient.discovery import build
from googleapiclient.http import MediaIoBaseUpload

load_dotenv()

GEMINI_MODEL = "gemini-2.5-flash"

IAC_BASE_URL = "https://server.iac.ac.il/api/v1/studentapi"
IAC_MAX_TOKENS = 30000

GOOGLE_SCOPES = [
    "https://www.googleapis.com/auth/gmail.modify",
    "https://www.googleapis.com/auth/calendar",
    "https://www.googleapis.com/auth/drive",
    "https://www.googleapis.com/auth/documents",
    "https://www.googleapis.com/auth/spreadsheets",
    "https://www.googleapis.com/auth/presentations",
]
CREDENTIALS_FILE = "credentials.json"
TOKEN_FILE = "gmail_token.json"

# ========================================================
# Page configuration
# ========================================================

st.set_page_config(page_title="הסוכן שלי", page_icon="🤖", layout="centered")

# עיצוב מותאם לבועות הצ'אט (מבוסס data-testid יציב של Streamlit)
st.markdown("""
<style>
/* בועת המשתמש - כחול gradient, מיושרת לימין */
.stChatMessage:has([data-testid="stChatMessageAvatarUser"]) {
    flex-direction: row-reverse;
    text-align: right;
}
.stChatMessage:has([data-testid="stChatMessageAvatarUser"]) [data-testid="stChatMessageContent"] {
    background: linear-gradient(135deg, #2563eb, #3b82f6);
    color: #ffffff;
    border-radius: 20px 20px 6px 20px;
    padding: 12px 16px;
    box-shadow: 0 4px 14px rgba(37, 99, 235, 0.25);
}
.stChatMessage:has([data-testid="stChatMessageAvatarUser"]) [data-testid="stChatMessageContent"] * {
    color: #ffffff !important;
}
/* בועת הסוכן - אפור בהיר, מיושרת לשמאל */
.stChatMessage:has([data-testid="stChatMessageAvatarAssistant"]) [data-testid="stChatMessageContent"] {
    background: #d8e0ea;
    color: #1e293b;
    border-radius: 20px 20px 20px 6px;
    padding: 12px 16px;
    box-shadow: 0 2px 8px rgba(0, 0, 0, 0.06);
}
/* אווטרים עגולים ומודגשים */
.stChatMessage [data-testid="stChatMessageAvatarUser"],
.stChatMessage [data-testid="stChatMessageAvatarAssistant"] {
    box-shadow: 0 2px 6px rgba(0, 0, 0, 0.12);
}
</style>
""", unsafe_allow_html=True)

st.title("🤖 הסוכן שלי")
st.caption("סוכן אחד עם כלים: קבצים, Word/PowerPoint, Gmail וחיפוש אינטרנט")

if "chat_history" not in st.session_state:
    st.session_state.chat_history = []

if "created_files" not in st.session_state:
    st.session_state.created_files = []

if "uploaded_paths" not in st.session_state:
    st.session_state.uploaded_paths = []

if "pending_email" not in st.session_state:
    st.session_state.pending_email = None

if "conversations" not in st.session_state:
    st.session_state.conversations = None  # None = טרם נטען מ-Drive

if "active_conversation_id" not in st.session_state:
    st.session_state.active_conversation_id = None

UPLOADS_DIR = Path("uploads")
UPLOADS_DIR.mkdir(exist_ok=True)


# ========================================================
# טוקנים
# ========================================================

def load_gemini_keys() -> list:
    """מחזיר את כל מפתחות ה-Gemini הזמינים (עד 3), לפי הסדר, בלי כפילויות."""
    keys = []
    for name in ("GEMINI_API_KEY", "GEMINI_API_KEY_2", "GEMINI_API_KEY_3",
                 "GEMINI_API_KEY_4", "GEMINI_API_KEY_5"):
        val = os.environ.get(name)
        if val and val not in keys:
            keys.append(val)
    return keys


def load_gemini_key() -> str:
    keys = load_gemini_keys()
    return keys[0] if keys else None


def load_iac_token() -> str:
    val = os.environ.get("IAC_TOKEN")
    if val and val != "sk-std-YOUR-KEY":
        return val

    env_path = Path(".env")
    if not env_path.exists():
        return None
    for line in env_path.read_text(encoding="utf-8").splitlines():
        if line.startswith("IAC_TOKEN="):
            val = line.split("=", 1)[1].strip().strip("'\"")
            if val and val != "sk-std-YOUR-KEY":
                return val
    return None


def is_gemini_quota_error(e: Exception) -> bool:
    return isinstance(e, genai_errors.ClientError) and "RESOURCE_EXHAUSTED" in str(e)


# ========================================================
# כלים: קבצים, shell, Gmail
# ========================================================

def get_google_creds():
    creds = None
    token_from_secret = os.environ.get("GMAIL_TOKEN_JSON")

    if Path(TOKEN_FILE).exists():
        creds = Credentials.from_authorized_user_file(TOKEN_FILE, GOOGLE_SCOPES)
    elif token_from_secret:
        creds = Credentials.from_authorized_user_info(json.loads(token_from_secret), GOOGLE_SCOPES)

    if not creds or not creds.valid:
        if creds and creds.expired and creds.refresh_token:
            creds.refresh(Request())
        elif token_from_secret:
            raise RuntimeError("טוקן ה-Google בענן פג ואין אפשרות לפתוח דפדפן לאישור מחדש בענן.")
        else:
            flow = InstalledAppFlow.from_client_secrets_file(CREDENTIALS_FILE, GOOGLE_SCOPES)
            creds = flow.run_local_server(port=0)
        if not token_from_secret:
            Path(TOKEN_FILE).write_text(creds.to_json(), encoding="utf-8")
    return creds


def get_gmail_service():
    return build("gmail", "v1", credentials=get_google_creds())


def get_calendar_service():
    return build("calendar", "v3", credentials=get_google_creds())


def get_drive_service():
    return build("drive", "v3", credentials=get_google_creds())


def get_docs_service():
    return build("docs", "v1", credentials=get_google_creds())


def get_sheets_service():
    return build("sheets", "v4", credentials=get_google_creds())


def get_slides_service():
    return build("slides", "v1", credentials=get_google_creds())


# ========================================================
# שמירת היסטוריית שיחות ב-Google Drive
# ========================================================

CONV_FILE_NAME = "agent_conversations.json"


def _drive_conv_file_id(service):
    resp = service.files().list(
        q=f"name='{CONV_FILE_NAME}' and trashed=false",
        spaces="drive", fields="files(id)",
    ).execute()
    files = resp.get("files", [])
    return files[0]["id"] if files else None


def load_conversations_from_drive() -> list:
    try:
        service = get_drive_service()
        fid = _drive_conv_file_id(service)
        if not fid:
            return []
        content = service.files().get_media(fileId=fid).execute()
        data = json.loads(content.decode("utf-8"))
        return data.get("conversations", [])
    except Exception:
        return []


def save_conversations_to_drive(conversations: list) -> bool:
    try:
        service = get_drive_service()
        payload = json.dumps({"conversations": conversations}, ensure_ascii=False).encode("utf-8")
        media = MediaIoBaseUpload(io.BytesIO(payload), mimetype="application/json")
        fid = _drive_conv_file_id(service)
        if fid:
            service.files().update(fileId=fid, media_body=media).execute()
        else:
            service.files().create(
                body={"name": CONV_FILE_NAME}, media_body=media, fields="id"
            ).execute()
        return True
    except Exception as e:
        st.toast(f"⚠️ שמירה ל-Drive נכשלה: {e}")
        return False


def _make_title(history: list) -> str:
    for m in history:
        if m["role"] == "user" and m.get("content"):
            first_line = m["content"].strip().split("\n")[0]
            return first_line[:40] or "שיחה חדשה"
    return "שיחה חדשה"


def start_new_conversation():
    st.session_state.active_conversation_id = None
    st.session_state.chat_history = []


def select_conversation(cid: str):
    conv = next((c for c in (st.session_state.conversations or []) if c["id"] == cid), None)
    if conv:
        st.session_state.active_conversation_id = cid
        st.session_state.chat_history = conv["messages"]


def persist_current_conversation():
    if not st.session_state.chat_history:
        return
    convs = st.session_state.conversations
    if convs is None:
        convs = []
        st.session_state.conversations = convs
    cid = st.session_state.active_conversation_id
    title = _make_title(st.session_state.chat_history)
    now = datetime.now().isoformat()
    if cid:
        conv = next((c for c in convs if c["id"] == cid), None)
        if conv:
            conv["messages"] = st.session_state.chat_history
            conv["title"] = title
            conv["updated"] = now
    else:
        cid = uuid.uuid4().hex
        convs.insert(0, {
            "id": cid, "title": title,
            "messages": st.session_state.chat_history, "updated": now,
        })
        st.session_state.active_conversation_id = cid
    save_conversations_to_drive(convs)


def delete_conversation(cid: str):
    convs = st.session_state.conversations or []
    st.session_state.conversations = [c for c in convs if c["id"] != cid]
    if st.session_state.active_conversation_id == cid:
        start_new_conversation()
    save_conversations_to_drive(st.session_state.conversations)


# ========================================================
# מעקב שימוש יומי: חיפושים וטוקנים
# ========================================================

USAGE_FILE = Path("usage_stats.json")


def _load_usage() -> dict:
    today = datetime.now().strftime("%Y-%m-%d")
    try:
        data = json.loads(USAGE_FILE.read_text(encoding="utf-8"))
    except Exception:
        data = {}
    if data.get("date") != today:
        data = {"date": today, "searches": 0, "tokens": 0, "requests": 0}
    return data


def _save_usage(data: dict):
    try:
        USAGE_FILE.write_text(json.dumps(data), encoding="utf-8")
    except Exception:
        pass


def track_search():
    data = _load_usage()
    data["searches"] = data.get("searches", 0) + 1
    _save_usage(data)


def track_tokens(count: int):
    if not count:
        return
    data = _load_usage()
    data["tokens"] = data.get("tokens", 0) + count
    data["requests"] = data.get("requests", 0) + 1
    _save_usage(data)


def tool_read_file(path: str) -> str:
    try:
        return Path(path).read_text(encoding="utf-8")
    except Exception as e:
        return f"שגיאה בקריאת הקובץ: {e}"


def tool_write_file(path: str, content: str) -> str:
    try:
        Path(path).write_text(content, encoding="utf-8")
        return f"הקובץ {path} נשמר בהצלחה."
    except Exception as e:
        return f"שגיאה בכתיבת הקובץ: {e}"


def tool_run_shell(command: str) -> str:
    import subprocess
    try:
        result = subprocess.run(command, shell=True, capture_output=True, text=True, timeout=60)
        output = result.stdout + result.stderr
        return output.strip() or "(הפקודה רצה בהצלחה, אין פלט)"
    except Exception as e:
        return f"שגיאה בהרצת הפקודה: {e}"


def tool_send_email(to: str, subject: str, body: str, attachments: list = None) -> str:
    st.session_state.pending_email = {
        "to": to, "subject": subject, "body": body, "attachments": attachments or [],
    }
    return (
        "המייל הוכן וממתין לאישור המשתמש בממשק (כפתורי אישור/ביטול). "
        "אל תקרא שוב לכלי הזה - רק הודע למשתמש שהמייל מוכן וממתין לאישורו."
    )


def _do_send_email(to: str, subject: str, body: str, attachments: list = None) -> str:
    try:
        service = get_gmail_service()

        if attachments:
            message = MIMEMultipart()
            message.attach(MIMEText(body))
            attached_names = []
            missing = []
            for path_str in attachments:
                p = Path(path_str)
                if not p.exists():
                    missing.append(path_str)
                    continue
                part = MIMEBase("application", "octet-stream")
                part.set_payload(p.read_bytes())
                encoders.encode_base64(part)
                part.add_header("Content-Disposition", f'attachment; filename="{p.name}"')
                message.attach(part)
                attached_names.append(p.name)
            if missing:
                return f"שגיאה: הקבצים הבאים לא נמצאו ולא נשלחו: {', '.join(missing)}"
        else:
            message = MIMEText(body)
            attached_names = []

        message["to"] = to
        message["subject"] = subject
        raw = base64.urlsafe_b64encode(message.as_bytes()).decode()
        service.users().messages().send(userId="me", body={"raw": raw}).execute()

        if attached_names:
            return f"המייל נשלח בהצלחה ל-{to} עם הקבצים המצורפים: {', '.join(attached_names)}."
        return f"המייל נשלח בהצלחה ל-{to}."
    except Exception as e:
        return f"שגיאה בשליחת המייל: {e}"


def tool_list_recent_emails(max_results: int = 5) -> str:
    try:
        service = get_gmail_service()
        results = service.users().messages().list(userId="me", maxResults=max_results).execute()
        messages = results.get("messages", [])
        if not messages:
            return "אין מיילים בתיבה."
        summaries = []
        for msg in messages:
            full = service.users().messages().get(
                userId="me", id=msg["id"], format="metadata",
                metadataHeaders=["From", "Subject"]
            ).execute()
            headers = {h["name"]: h["value"] for h in full["payload"]["headers"]}
            snippet = full.get("snippet", "")
            summaries.append(f"מאת: {headers.get('From', '?')} | נושא: {headers.get('Subject', '?')} | תקציר: {snippet}")
        return "\n".join(summaries)
    except Exception as e:
        return f"שגיאה בקריאת המיילים: {e}"


# ========================================================
# כלים: Google Calendar
# ========================================================

def tool_create_calendar_event(summary: str, start: str, end: str, description: str = "") -> str:
    try:
        service = get_calendar_service()
        event = {
            "summary": summary,
            "description": description,
            "start": {"dateTime": start, "timeZone": "Asia/Jerusalem"},
            "end": {"dateTime": end, "timeZone": "Asia/Jerusalem"},
        }
        created = service.events().insert(calendarId="primary", body=event).execute()
        return f"האירוע '{summary}' נוצר בהצלחה בלוח השנה. קישור: {created.get('htmlLink', '')}"
    except Exception as e:
        return f"שגיאה ביצירת אירוע בלוח השנה: {e}"


def tool_list_calendar_events(max_results: int = 10) -> str:
    try:
        from datetime import datetime, timezone
        service = get_calendar_service()
        now = datetime.now(timezone.utc).isoformat()
        results = service.events().list(
            calendarId="primary", timeMin=now, maxResults=max_results,
            singleEvents=True, orderBy="startTime",
        ).execute()
        events = results.get("items", [])
        if not events:
            return "אין אירועים קרובים בלוח השנה."
        summaries = []
        for ev in events:
            start = ev["start"].get("dateTime", ev["start"].get("date"))
            summaries.append(f"{start} | {ev.get('summary', '(ללא כותרת)')}")
        return "\n".join(summaries)
    except Exception as e:
        return f"שגיאה בקריאת לוח השנה: {e}"


# ========================================================
# כלים: Google Drive / Docs / Sheets / Slides
# ========================================================

def tool_search_drive_files(query: str, max_results: int = 10) -> str:
    try:
        service = get_drive_service()
        results = service.files().list(
            q=f"name contains '{query}'", pageSize=max_results,
            fields="files(id, name, mimeType, webViewLink)",
        ).execute()
        files = results.get("files", [])
        if not files:
            return "לא נמצאו קבצים מתאימים ב-Drive."
        return "\n".join(f"{f['name']} ({f['mimeType']}) | {f.get('webViewLink', '')}" for f in files)
    except Exception as e:
        return f"שגיאה בחיפוש ב-Drive: {e}"


def tool_create_google_doc(title: str, content: str) -> str:
    try:
        docs_service = get_docs_service()
        doc = docs_service.documents().create(body={"title": title}).execute()
        doc_id = doc["documentId"]
        if content:
            docs_service.documents().batchUpdate(
                documentId=doc_id,
                body={"requests": [{"insertText": {"location": {"index": 1}, "text": content}}]},
            ).execute()
        return f"מסמך Google Docs '{title}' נוצר בהצלחה: https://docs.google.com/document/d/{doc_id}/edit"
    except Exception as e:
        return f"שגיאה ביצירת מסמך Google Docs: {e}"


def tool_create_google_sheet(title: str, headers: list, rows: list) -> str:
    try:
        sheets_service = get_sheets_service()
        sheet = sheets_service.spreadsheets().create(body={"properties": {"title": title}}).execute()
        sheet_id = sheet["spreadsheetId"]
        values = ([headers] if headers else []) + (rows or [])
        if values:
            sheets_service.spreadsheets().values().update(
                spreadsheetId=sheet_id, range="A1",
                valueInputOption="RAW", body={"values": values},
            ).execute()
        return f"גיליון Google Sheets '{title}' נוצר בהצלחה: https://docs.google.com/spreadsheets/d/{sheet_id}/edit"
    except Exception as e:
        return f"שגיאה ביצירת גיליון Google Sheets: {e}"


def tool_create_google_slides(title: str, slides: list) -> str:
    try:
        slides_service = get_slides_service()
        presentation = slides_service.presentations().create(body={"title": title}).execute()
        pres_id = presentation["presentationId"]

        requests_batch = []
        for slide_data in slides:
            slide_id = f"slide_{len(requests_batch)}"
            requests_batch.append({
                "createSlide": {
                    "objectId": slide_id,
                    "slideLayoutReference": {"predefinedLayout": "TITLE_AND_BODY"},
                    "placeholderIdMappings": [
                        {"layoutPlaceholder": {"type": "TITLE"}, "objectId": f"{slide_id}_title"},
                        {"layoutPlaceholder": {"type": "BODY"}, "objectId": f"{slide_id}_body"},
                    ],
                }
            })
        if requests_batch:
            slides_service.presentations().batchUpdate(
                presentationId=pres_id, body={"requests": requests_batch}
            ).execute()

        text_requests = []
        for i, slide_data in enumerate(slides):
            slide_id = f"slide_{i}"
            text_requests.append({
                "insertText": {"objectId": f"{slide_id}_title", "text": slide_data.get("title", "")}
            })
            text_requests.append({
                "insertText": {"objectId": f"{slide_id}_body", "text": slide_data.get("content", "")}
            })
        if text_requests:
            slides_service.presentations().batchUpdate(
                presentationId=pres_id, body={"requests": text_requests}
            ).execute()

        return f"מצגת Google Slides '{title}' נוצרה בהצלחה: https://docs.google.com/presentation/d/{pres_id}/edit"
    except Exception as e:
        return f"שגיאה ביצירת מצגת Google Slides: {e}"


# ========================================================
# כלים: Word / PowerPoint
# ========================================================

def _register_created_file(path: str):
    if path not in st.session_state.created_files:
        st.session_state.created_files.append(path)


def tool_create_word_doc(path: str, title: str, paragraphs: list) -> str:
    try:
        if not path.endswith(".docx"):
            path += ".docx"
        doc = Document()
        if title:
            doc.add_heading(title, level=1)
        for p in paragraphs:
            doc.add_paragraph(p)
        doc.save(path)
        _register_created_file(path)
        return f"קובץ Word נוצר בהצלחה: {path}"
    except Exception as e:
        return f"שגיאה ביצירת קובץ Word: {e}"


def tool_create_pptx(path: str, slides: list) -> str:
    try:
        if not path.endswith(".pptx"):
            path += ".pptx"
        prs = Presentation()
        title_layout = prs.slide_layouts[1]
        for slide_data in slides:
            slide = prs.slides.add_slide(title_layout)
            slide.shapes.title.text = slide_data.get("title", "")
            body = slide.placeholders[1].text_frame
            content = slide_data.get("content", "")
            lines = content.split("\n") if isinstance(content, str) else content
            for i, line in enumerate(lines):
                if i == 0:
                    body.text = line
                else:
                    body.add_paragraph().text = line
        prs.save(path)
        _register_created_file(path)
        return f"מצגת PowerPoint נוצרה בהצלחה: {path}"
    except Exception as e:
        return f"שגיאה ביצירת מצגת: {e}"


def tool_create_excel(path: str, sheet_name: str, headers: list, rows: list) -> str:
    try:
        if not path.endswith(".xlsx"):
            path += ".xlsx"
        wb = Workbook()
        ws = wb.active
        ws.title = sheet_name or "Sheet1"
        if headers:
            ws.append(headers)
        for row in rows:
            ws.append(row)
        wb.save(path)
        _register_created_file(path)
        return f"קובץ Excel נוצר בהצלחה: {path}"
    except Exception as e:
        return f"שגיאה ביצירת קובץ Excel: {e}"


# ========================================================
# כלי: חיפוש אינטרנט (Google Custom Search - חינמי, עצמאי)
# ========================================================

GOOGLE_SEARCH_URL = "https://www.googleapis.com/customsearch/v1"


def _search_google(query: str) -> str:
    api_key = os.environ.get("GOOGLE_SEARCH_API_KEY")
    engine_id = os.environ.get("GOOGLE_SEARCH_ENGINE_ID")
    if not api_key or not engine_id:
        raise RuntimeError("חסר GOOGLE_SEARCH_API_KEY או GOOGLE_SEARCH_ENGINE_ID")
    params = {"key": api_key, "cx": engine_id, "q": query, "num": 5, "hl": "he"}
    r = requests.get(GOOGLE_SEARCH_URL, params=params, timeout=30)
    if r.status_code == 429:
        raise RuntimeError("מגבלת חיפושים יומית של Google (100 ביום)")
    r.raise_for_status()
    items = r.json().get("items", [])
    if not items:
        return "לא נמצאו תוצאות."
    return "\n\n".join(
        f"• {it.get('title', '')}\n  {it.get('snippet', '')}\n  {it.get('link', '')}"
        for it in items
    )


def _search_duckduckgo(query: str) -> str:
    # גיבוי חינמי בלי מפתח - DuckDuckGo HTML
    r = requests.get(
        "https://html.duckduckgo.com/html/",
        params={"q": query},
        headers={"User-Agent": "Mozilla/5.0"},
        timeout=30,
    )
    r.raise_for_status()
    import re
    from html import unescape
    results = []
    blocks = re.findall(
        r'<a[^>]*class="result__a"[^>]*href="([^"]*)"[^>]*>(.*?)</a>.*?'
        r'<a[^>]*class="result__snippet"[^>]*>(.*?)</a>',
        r.text, re.DOTALL,
    )
    for link, title, snippet in blocks[:5]:
        clean = lambda s: unescape(re.sub(r"<[^>]+>", "", s)).strip()
        results.append(f"• {clean(title)}\n  {clean(snippet)}\n  {link}")
    return "\n\n".join(results) if results else "לא נמצאו תוצאות."


def tool_web_search(query: str) -> str:
    track_search()
    try:
        return _search_google(query)
    except Exception:
        try:
            st.toast("🦆 Google לא זמין - מחפש דרך DuckDuckGo")
            return _search_duckduckgo(query)
        except Exception as e:
            return f"שגיאה בחיפוש אינטרנט: {e}"


def call_tool(name: str, args: dict) -> str:
    func = TOOL_FUNCTIONS.get(name)
    if not func:
        return f"כלי לא מוכר: {name}"
    try:
        return func(**args)
    except TypeError:
        if set(args.keys()) - {"parameters", "arguments", "args"} == set() and len(args) == 1:
            inner = next(iter(args.values()))
            if isinstance(inner, dict):
                return func(**inner)
        raise


TOOL_FUNCTIONS = {
    "read_file": tool_read_file,
    "write_file": tool_write_file,
    "run_shell": tool_run_shell,
    "send_email": tool_send_email,
    "list_recent_emails": tool_list_recent_emails,
    "create_calendar_event": tool_create_calendar_event,
    "list_calendar_events": tool_list_calendar_events,
    "search_drive_files": tool_search_drive_files,
    "create_google_doc": tool_create_google_doc,
    "create_google_sheet": tool_create_google_sheet,
    "create_google_slides": tool_create_google_slides,
    "create_word_doc": tool_create_word_doc,
    "create_pptx": tool_create_pptx,
    "create_excel": tool_create_excel,
    "web_search": tool_web_search,
}

TOOLS_SCHEMA = [
    {
        "name": "read_file",
        "description": "קורא את תוכן הקובץ הנתון ומחזיר אותו כטקסט.",
        "parameters": {
            "type": "object",
            "properties": {"path": {"type": "string", "description": "נתיב לקובץ לקריאה"}},
            "required": ["path"],
        },
    },
    {
        "name": "write_file",
        "description": "כותב תוכן לקובץ נתון (יוצר אותו אם לא קיים, או דורס אותו אם קיים).",
        "parameters": {
            "type": "object",
            "properties": {
                "path": {"type": "string", "description": "נתיב לקובץ לכתיבה"},
                "content": {"type": "string", "description": "התוכן לכתוב לקובץ"},
            },
            "required": ["path", "content"],
        },
    },
    {
        "name": "run_shell",
        "description": "מריץ פקודת shell ומחזיר את הפלט שלה (stdout/stderr).",
        "parameters": {
            "type": "object",
            "properties": {"command": {"type": "string", "description": "פקודת shell להרצה"}},
            "required": ["command"],
        },
    },
    {
        "name": "send_email",
        "description": "שולח מייל מתיבת הדואר של המשתמש (Gmail) לכתובת נתונה. ניתן לצרף קבצים (Word/PowerPoint/Excel/תמונות) ולהוסיף קישורים בתוך גוף הטקסט.",
        "parameters": {
            "type": "object",
            "properties": {
                "to": {"type": "string", "description": "כתובת המייל של הנמען"},
                "subject": {"type": "string", "description": "נושא המייל"},
                "body": {"type": "string", "description": "תוכן המייל (אפשר לכלול קישורים כטקסט רגיל, הם יוצגו כניתנים ללחיצה)"},
                "attachments": {
                    "type": "array",
                    "items": {"type": "string"},
                    "description": "רשימת נתיבי קבצים לצירוף למייל (למשל קובץ Word/PowerPoint/Excel שנוצר קודם, או קובץ תמונה שהועלה)",
                },
            },
            "required": ["to", "subject", "body"],
        },
    },
    {
        "name": "list_recent_emails",
        "description": "מחזיר רשימה של המיילים האחרונים בתיבת הדואר (נושא, שולח, תקציר).",
        "parameters": {
            "type": "object",
            "properties": {
                "max_results": {"type": "integer", "description": "כמה מיילים להחזיר (ברירת מחדל 5)"}
            },
            "required": [],
        },
    },
    {
        "name": "create_calendar_event",
        "description": "יוצר אירוע אמיתי בלוח השנה של Google של המשתמש (תזכורת, פגישה, וכו').",
        "parameters": {
            "type": "object",
            "properties": {
                "summary": {"type": "string", "description": "כותרת/שם האירוע"},
                "start": {"type": "string", "description": "זמן התחלה בפורמט ISO 8601, למשל '2026-06-26T15:00:00'"},
                "end": {"type": "string", "description": "זמן סיום בפורמט ISO 8601, למשל '2026-06-26T16:00:00'"},
                "description": {"type": "string", "description": "תיאור נוסף לאירוע (אופציונלי)"},
            },
            "required": ["summary", "start", "end"],
        },
    },
    {
        "name": "list_calendar_events",
        "description": "מחזיר רשימה של האירועים הקרובים בלוח השנה של המשתמש.",
        "parameters": {
            "type": "object",
            "properties": {
                "max_results": {"type": "integer", "description": "כמה אירועים להחזיר (ברירת מחדל 10)"}
            },
            "required": [],
        },
    },
    {
        "name": "search_drive_files",
        "description": "מחפש קבצים בגוגל דרייב של המשתמש לפי שם.",
        "parameters": {
            "type": "object",
            "properties": {
                "query": {"type": "string", "description": "מילת חיפוש (חלק מהשם)"},
                "max_results": {"type": "integer", "description": "כמה תוצאות להחזיר (ברירת מחדל 10)"},
            },
            "required": ["query"],
        },
    },
    {
        "name": "create_google_doc",
        "description": "יוצר מסמך Google Docs אמיתי בדרייב של המשתמש (לא קובץ מקומי).",
        "parameters": {
            "type": "object",
            "properties": {
                "title": {"type": "string", "description": "כותרת המסמך"},
                "content": {"type": "string", "description": "תוכן טקסטואלי להכניס למסמך"},
            },
            "required": ["title"],
        },
    },
    {
        "name": "create_google_sheet",
        "description": "יוצר גיליון Google Sheets אמיתי בדרייב של המשתמש (ה'Excel' של גוגל, לא קובץ מקומי).",
        "parameters": {
            "type": "object",
            "properties": {
                "title": {"type": "string", "description": "כותרת הגיליון"},
                "headers": {
                    "type": "array", "items": {"type": "string"},
                    "description": "כותרות העמודות (שורה ראשונה)",
                },
                "rows": {
                    "type": "array", "items": {"type": "array", "items": {}},
                    "description": "רשימת שורות נתונים",
                },
            },
            "required": ["title", "rows"],
        },
    },
    {
        "name": "create_google_slides",
        "description": "יוצר מצגת Google Slides אמיתית בדרייב של המשתמש (לא קובץ מקומי).",
        "parameters": {
            "type": "object",
            "properties": {
                "title": {"type": "string", "description": "כותרת המצגת"},
                "slides": {
                    "type": "array",
                    "items": {
                        "type": "object",
                        "properties": {
                            "title": {"type": "string", "description": "כותרת השקופית"},
                            "content": {"type": "string", "description": "תוכן השקופית"},
                        },
                    },
                    "description": "רשימת שקופיות",
                },
            },
            "required": ["title", "slides"],
        },
    },
    {
        "name": "create_word_doc",
        "description": "יוצר קובץ Word (.docx) עם כותרת ופסקאות טקסט.",
        "parameters": {
            "type": "object",
            "properties": {
                "path": {"type": "string", "description": "נתיב/שם הקובץ לשמירה, למשל 'document.docx'"},
                "title": {"type": "string", "description": "כותרת המסמך"},
                "paragraphs": {
                    "type": "array",
                    "items": {"type": "string"},
                    "description": "רשימת פסקאות טקסט להוספה למסמך",
                },
            },
            "required": ["path", "paragraphs"],
        },
    },
    {
        "name": "create_pptx",
        "description": "יוצר מצגת PowerPoint (.pptx) עם רשימת שקופיות, כל שקופית עם כותרת ותוכן.",
        "parameters": {
            "type": "object",
            "properties": {
                "path": {"type": "string", "description": "נתיב/שם הקובץ לשמירה, למשל 'presentation.pptx'"},
                "slides": {
                    "type": "array",
                    "items": {
                        "type": "object",
                        "properties": {
                            "title": {"type": "string", "description": "כותרת השקופית"},
                            "content": {"type": "string", "description": "תוכן השקופית (אפשר כמה שורות מופרדות ב-\\n)"},
                        },
                    },
                    "description": "רשימת שקופיות, כל אחת עם title ו-content",
                },
            },
            "required": ["path", "slides"],
        },
    },
    {
        "name": "create_excel",
        "description": "יוצר קובץ Excel (.xlsx) עם כותרות עמודות ושורות נתונים.",
        "parameters": {
            "type": "object",
            "properties": {
                "path": {"type": "string", "description": "נתיב/שם הקובץ לשמירה, למשל 'data.xlsx'"},
                "sheet_name": {"type": "string", "description": "שם הגיליון"},
                "headers": {
                    "type": "array",
                    "items": {"type": "string"},
                    "description": "כותרות העמודות (שורה ראשונה)",
                },
                "rows": {
                    "type": "array",
                    "items": {"type": "array", "items": {}},
                    "description": "רשימת שורות, כל שורה היא רשימת ערכים",
                },
            },
            "required": ["path", "rows"],
        },
    },
    {
        "name": "web_search",
        "description": "מחפש מידע עדכני באינטרנט ומחזיר תקציר תשובה. שימושי לשאלות על אירועים עדכניים, מחירים, חדשות וכו'.",
        "parameters": {
            "type": "object",
            "properties": {"query": {"type": "string", "description": "שאילתת החיפוש"}},
            "required": ["query"],
        },
    },
]

MAX_TOOL_ROUNDS = 10

SYSTEM_INSTRUCTION = (
    "אתה סוכן AI פעיל עם גישה אמיתית לכלים הבאים, ואתה חייב להשתמש בהם בפועל ולא רק לתאר מה "
    "אפשר לעשות:\n"
    "- send_email: שולח מייל אמיתי מתיבת ה-Gmail המחוברת של המשתמש. כשמתבקש לשלוח מייל, "
    "תקרא לכלי הזה ממש - אל תכתוב שאינך יכול לשלוח מיילים. אם המשתמש מבקש לצרף קובץ "
    "(Word/PowerPoint/Excel/תמונה שנוצר או הועלה קודם), העבר את הנתיב שלו בפרמטר attachments. "
    "קישורים אפשר לכתוב בתוך גוף ה-body כטקסט רגיל.\n"
    "- list_recent_emails: קורא מיילים אמיתיים מהתיבה.\n"
    "- create_calendar_event, list_calendar_events: יוצרים/קוראים אירועים אמיתיים בלוח השנה "
    "של Google של המשתמש. כשמתבקש לקבוע/לתזכר משהו בתאריך ושעה, תקרא לכלי הזה ממש.\n"
    "- search_drive_files: מחפש קבצים אמיתיים בגוגל דרייב של המשתמש.\n"
    "- create_google_doc, create_google_sheet, create_google_slides: יוצרים מסמך/גיליון/מצגת "
    "אמיתיים בגוגל דרייב של המשתמש (לא קובץ מקומי). אם המשתמש מבקש Google Docs/Sheets/Slides "
    "באופן מפורש, או 'Excel של גוגל', תשתמש בכלים האלה ולא בקבצים המקומיים.\n"
    "- create_word_doc, create_pptx, create_excel: יוצרים קבצים מקומיים (.docx/.pptx/.xlsx) "
    "בדיסק שהמשתמש יכול להוריד.\n"
    "- web_search: מחפש מידע עדכני באינטרנט באמת.\n"
    "- read_file, write_file, run_shell: גישה אמיתית למערכת הקבצים ולמסוף.\n"
    "לעולם אל תגיד שאינך מסוגל לבצע פעולה שיש לך כלי בשבילה - תמיד תקרא לכלי המתאים ותבצע "
    "את הבקשה בפועל.\n"
    "חשוב מאוד: ענה תמיד באותה שפה שבה המשתמש פנה אליך בהודעה. אם כתב בעברית - ענה בעברית; "
    "אם כתב באנגלית - ענה באנגלית; וכן הלאה לכל שפה.\n"
    "אל תציג למשתמש את החשיבה הפנימית שלך. לעולם אל תכתוב 'THOUGHT:' או מחשבות/הסברים על "
    "מה שאתה עומד לעשות - החזר אך ורק את התשובה הסופית והנקייה למשתמש."
)


# ========================================================
# Gemini
# ========================================================

def _strip_thoughts(text: str) -> str:
    """מסיר חשיבה פנימית שדולפת (THOUGHT או פסקת מחשבה באנגלית) ומשאיר רק את התשובה הנקייה."""
    if not text:
        return text
    import re
    # 1) מסיר בלוק THOUGHT/THINKING מפורש בתחילת הטקסט
    cleaned = re.sub(r'(?is)^\s*(THOUGHT|THINKING|מחשבה)\s*:.*?(?=\n|[א-ת])', '', text, count=1)

    # 2) רשת ביטחון: אם יש פסקת "חשיבה" באנגלית לפני התשובה בעברית -
    #    (מזוהה לפי ניסוחי מטא כמו "The user", "I should", "This is") - חותכים עד לעברית הראשונה
    thought_markers = ("the user", "i should", "i will", "this is great",
                       "this implies", "i need to", "let me", "i can")
    head = cleaned[:400].lower()
    if any(m in head for m in thought_markers):
        heb = re.search(r'[א-ת]', cleaned)
        if heb and heb.start() > 0:
            cleaned = cleaned[heb.start():]

    return cleaned.strip() or text.strip()


def get_gemini_client(api_key: str = None):
    api_key = api_key or load_gemini_key()
    if not api_key:
        raise RuntimeError("לא נמצא GEMINI_API_KEY בקובץ .env")
    return genai.Client(api_key=api_key)


def generate_with_retry(client, **kwargs):
    for attempt in range(3):
        try:
            return client.models.generate_content(**kwargs)
        except genai_errors.ServerError:
            time.sleep(3 * (attempt + 1))
    raise RuntimeError("Gemini לא הגיב (עומס שרת חזק).")


# פורמטים של תמונה ש-Gemini תומך בהם ישירות
GEMINI_IMAGE_MIMES = {"image/png", "image/jpeg", "image/webp", "image/heic", "image/heif"}


def _to_supported_image(p: Path, mime: str):
    """ממיר תמונות בפורמט לא נתמך (TIFF/BMP/GIF וכו') ל-PNG כדי ש-Gemini יוכל לראות אותן."""
    if mime in GEMINI_IMAGE_MIMES:
        return p.read_bytes(), mime
    try:
        from PIL import Image
        import io as _io
        img = Image.open(p)
        if img.mode not in ("RGB", "RGBA"):
            img = img.convert("RGB")
        buf = _io.BytesIO()
        img.save(buf, format="PNG")
        return buf.getvalue(), "image/png"
    except Exception:
        # אם ההמרה נכשלה - שולחים כמו שהוא ומקווים לטוב
        return p.read_bytes(), mime


def _extract_office_text(p: Path) -> str:
    """מחלץ טקסט מקבצי Office (Word/PowerPoint/Excel) כי Gemini לא קורא אותם ישירות."""
    ext = p.suffix.lower()
    try:
        if ext == ".docx":
            doc = Document(str(p))
            return "\n".join(par.text for par in doc.paragraphs if par.text.strip())
        if ext == ".pptx":
            prs = Presentation(str(p))
            lines = []
            for i, slide in enumerate(prs.slides, 1):
                lines.append(f"--- שקופית {i} ---")
                for shape in slide.shapes:
                    if shape.has_text_frame and shape.text_frame.text.strip():
                        lines.append(shape.text_frame.text)
            return "\n".join(lines)
        if ext == ".xlsx":
            wb = load_workbook(str(p), read_only=True, data_only=True)
            lines = []
            for ws in wb.worksheets:
                lines.append(f"--- גיליון {ws.title} ---")
                for row in ws.iter_rows(values_only=True):
                    cells = [str(c) for c in row if c is not None]
                    if cells:
                        lines.append(" | ".join(cells))
            return "\n".join(lines)
    except Exception as e:
        return f"(לא הצלחתי לחלץ טקסט מהקובץ {p.name}: {e})"
    return ""


def _build_media_parts(client, paths: list) -> list:
    """יוצר חלקי מדיה (תמונה/סרטון/PDF/Office) כדי ש-Gemini יוכל לראות ולנתח את התוכן."""
    import mimetypes
    office_exts = {".docx", ".pptx", ".xlsx"}
    parts = []
    for path_str in paths or []:
        p = Path(path_str)
        if not p.exists():
            continue
        ext = p.suffix.lower()
        mime, _ = mimetypes.guess_type(str(p))
        mime = mime or "application/octet-stream"
        try:
            # קבצי Office - מחלצים טקסט ומעבירים כטקסט
            if ext in office_exts:
                text = _extract_office_text(p)
                if text:
                    parts.append(types.Part(text=f"\n\nתוכן הקובץ '{p.name}':\n{text}"))
            # תמונות קטנות - שולחים ישירות (inline), עם המרה לפורמט נתמך אם צריך
            elif mime.startswith("image/") and p.stat().st_size < 15 * 1024 * 1024:
                data, send_mime = _to_supported_image(p, mime)
                parts.append(types.Part.from_bytes(data=data, mime_type=send_mime))
            else:
                # סרטונים/PDF/קבצים גדולים - מעלים דרך Files API
                uploaded = client.files.upload(file=str(p))
                parts.append(types.Part.from_uri(file_uri=uploaded.uri, mime_type=uploaded.mime_type))
        except Exception as e:
            st.toast(f"⚠️ לא הצלחתי לצרף את {p.name}: {e}")
    return parts


def call_agent_gemini(history: list, api_key: str = None, attachments: list = None) -> str:
    client = get_gemini_client(api_key)
    tool_config = types.Tool(function_declarations=[
        types.FunctionDeclaration(name=t["name"], description=t["description"], parameters=t["parameters"])
        for t in TOOLS_SCHEMA
    ])
    contents = [
        types.Content(role="model" if m["role"] == "assistant" else "user", parts=[types.Part(text=m["content"])])
        for m in history if m["role"] in ("user", "assistant") and m.get("content")
    ]

    # מצרף את התמונות/סרטונים עצמם להודעה האחרונה כדי שהסוכן יראה את התוכן
    if attachments and contents and contents[-1].role == "user":
        media_parts = _build_media_parts(client, attachments)
        contents[-1].parts.extend(media_parts)

    for _ in range(MAX_TOOL_ROUNDS):
        response = generate_with_retry(
            client, model=GEMINI_MODEL, contents=contents,
            config=types.GenerateContentConfig(
                tools=[tool_config], system_instruction=SYSTEM_INSTRUCTION
            ),
        )
        candidate = response.candidates[0]
        contents.append(candidate.content)

        usage = getattr(response, "usage_metadata", None)
        if usage:
            track_tokens(getattr(usage, "total_token_count", 0) or 0)

        function_calls = [p.function_call for p in candidate.content.parts if p.function_call]
        if not function_calls:
            # מסננים החוצה חלקי "חשיבה" (thought) שהמודל מחזיר, ומשאירים רק את התשובה
            answer = "".join(
                p.text for p in candidate.content.parts
                if p.text and not getattr(p, "thought", False)
            )
            return _strip_thoughts(answer)

        response_parts = []
        for call in function_calls:
            st.toast(f"🔧 קורא לכלי: {call.name}")
            result = call_tool(call.name, dict(call.args))
            response_parts.append(types.Part.from_function_response(name=call.name, response={"result": result}))
        contents.append(types.Content(role="user", parts=response_parts))

    return "הגעתי למספר המקסימלי של שלבים מבלי לסיים את המשימה."


# ========================================================
# IAC (גיבוי)
# ========================================================

def _iac_tools_schema() -> list:
    return [{"type": "function", "function": t} for t in TOOLS_SCHEMA]


def call_agent_iac(history: list) -> str:
    token = load_iac_token()
    if not token:
        raise RuntimeError("אין גם טוקן IAC לגיבוי.")
    headers = {"Authorization": f"Bearer {token}"}
    messages = [{"role": "system", "content": SYSTEM_INSTRUCTION}] + [
        {"role": m["role"], "content": m["content"]} for m in history if m.get("content")
    ]

    for _ in range(MAX_TOOL_ROUNDS):
        payload = {"messages": messages, "tools": _iac_tools_schema(), "max_completion_tokens": IAC_MAX_TOKENS}

        for attempt in range(3):
            try:
                r = requests.post(f"{IAC_BASE_URL}/chat/completions", headers=headers, json=payload, timeout=120)
                break
            except requests.exceptions.ReadTimeout:
                if attempt == 2:
                    raise RuntimeError("השרת של המכללה לא הגיב בזמן (timeout חזרתי). נסה שוב בעוד דקה.")
                time.sleep(3)

        r.raise_for_status()
        data = r.json()
        if "choices" not in data:
            raise RuntimeError(data.get("details") or data.get("error") or str(data))
        choice = data["choices"][0]
        message = choice["message"]
        tool_calls = message.get("tool_calls")
        if not tool_calls:
            content = message.get("content")
            if not content and choice.get("finish_reason") == "length":
                raise RuntimeError(
                    "השרת הגיע למגבלת הטוקנים בלי לסיים תשובה (כל התקציב הלך על חשיבה פנימית). נסה לשאול שאלה קצרה יותר."
                )
            return _strip_thoughts(content) if content else str(data)
        messages.append(message)
        for call in tool_calls:
            name = call["function"]["name"]
            args = json.loads(call["function"]["arguments"] or "{}")
            st.toast(f"🔧 קורא לכלי (גיבוי): {name}")
            result = call_tool(name, args)
            messages.append({"role": "tool", "tool_call_id": call["id"], "content": result})

    return "הגעתי למספר המקסימלי של שלבים מבלי לסיים את המשימה."


def call_agent(history: list, attachments: list = None) -> str:
    keys = load_gemini_keys()
    # מנסה כל מפתח Gemini בתורו; אם אחד מיצה את המכסה - עובר לבא
    for i, key in enumerate(keys):
        try:
            return call_agent_gemini(history, api_key=key, attachments=attachments)
        except Exception as e:
            if is_gemini_quota_error(e):
                if i < len(keys) - 1:
                    st.toast(f"⚠️ מפתח Gemini {i + 1} מיצה מכסה - עובר למפתח {i + 2}")
                    continue
                # כל מפתחות Gemini מוצו. אם יש תמונה/מדיה - IAC לא יודע לראות אותה,
                # אז עדיף להגיד למשתמש את האמת במקום תשובה שגויה מ-IAC.
                if attachments:
                    return (
                        "⚠️ המכסה היומית של כל מפתחות Gemini נגמרה, ורק Gemini יודע לראות תמונות "
                        "(הגיבוי IAC מעבד טקסט בלבד). לכן אני לא יכול לנתח את התמונה כרגע. "
                        "נסה שוב מאוחר יותר (המכסה מתאפסת כל יום), או שלח לי את הטקסט שבתמונה כטקסט."
                    )
                if load_iac_token():
                    st.toast("⚠️ כל מפתחות Gemini מוצו - עובר לגיבוי IAC")
                    return call_agent_iac(history)
            raise
    # אין אף מפתח Gemini - מנסים IAC ישירות
    if load_iac_token():
        return call_agent_iac(history)
    raise RuntimeError("לא נמצא אף מפתח Gemini ואין טוקן IAC לגיבוי.")


# ========================================================
# Sidebar
# ========================================================

with st.sidebar:
    st.header("הגדרות")
    st.caption("סוכן אחד עם כל הכלים: קבצים, Word/PowerPoint/Excel, Gmail, חיפוש אינטרנט")

    st.divider()
    st.subheader("📊 שימוש היום")
    usage = _load_usage()
    col_a, col_b = st.columns(2)
    col_a.metric("🔍 חיפושים", f"{usage.get('searches', 0)} / 100")
    col_b.metric("🧠 טוקנים", f"{usage.get('tokens', 0):,}")
    st.caption(f"בקשות ל-Gemini היום: {usage.get('requests', 0)} | מתאפס כל יום")

    st.divider()
    st.subheader("💬 השיחות שלי")

    # טעינה חד-פעמית של השיחות השמורות מ-Google Drive
    if st.session_state.conversations is None and load_gemini_key():
        with st.spinner("טוען שיחות מ-Google Drive..."):
            st.session_state.conversations = load_conversations_from_drive()

    if st.button("➕ שיחה חדשה", use_container_width=True):
        start_new_conversation()
        st.rerun()

    for conv in st.session_state.conversations or []:
        is_active = conv["id"] == st.session_state.active_conversation_id
        col_open, col_del = st.columns([5, 1])
        label = ("🟢 " if is_active else "") + conv["title"]
        if col_open.button(label, key=f"conv_{conv['id']}", use_container_width=True):
            select_conversation(conv["id"])
            st.rerun()
        if col_del.button("🗑", key=f"del_{conv['id']}"):
            delete_conversation(conv["id"])
            st.rerun()

    st.divider()
    st.subheader("📎 צרף קובץ")
    uploaded_file = st.file_uploader(
        "תמונה, סרטון או מסמך", type=None, key="file_upload_widget"
    )
    if uploaded_file is not None:
        save_path = UPLOADS_DIR / uploaded_file.name
        save_path.write_bytes(uploaded_file.getbuffer())
        path_str = str(save_path)
        if path_str not in st.session_state.uploaded_paths:
            st.session_state.uploaded_paths.append(path_str)
            st.success(f"הקובץ הועלה: {uploaded_file.name}")

    if st.session_state.uploaded_paths:
        st.caption("קבצים מצורפים לשיחה:")
        for p in st.session_state.uploaded_paths:
            st.text(Path(p).name)
        if st.button("נקה קבצים מצורפים"):
            st.session_state.uploaded_paths = []
            st.rerun()

    if st.session_state.created_files:
        st.divider()
        st.subheader("⬇️ קבצים שנוצרו")
        for fpath in st.session_state.created_files:
            p = Path(fpath)
            if p.exists():
                st.download_button(
                    f"הורד {p.name}", data=p.read_bytes(), file_name=p.name, key=f"dl_{fpath}"
                )


# ========================================================
# Token check
# ========================================================

if not load_gemini_key():
    st.error("לא נמצא מפתח תקין בקובץ .env")
    st.info("יש להוסיף לקובץ .env שורה כזו:")
    st.code("GEMINI_API_KEY=your_key_here")
    st.stop()


# ========================================================
# Display history
# ========================================================

if not st.session_state.chat_history:
    with st.chat_message("assistant"):
        st.markdown("היי! אני יכול לעזור עם קבצים, Word/PowerPoint, מיילים וחיפוש באינטרנט 🙂")

for msg in st.session_state.chat_history:
    if msg["role"] not in ("user", "assistant"):
        continue
    content = msg.get("content")
    if not content:
        continue
    with st.chat_message(msg["role"]):
        st.markdown(content)


# ========================================================
# אישור שליחת מייל
# ========================================================

if st.session_state.pending_email:
    pe = st.session_state.pending_email
    with st.chat_message("assistant"):
        st.markdown("📧 **הסוכן רוצה לשלוח מייל:**")
        st.markdown(f"**אל:** {pe['to']}\n\n**נושא:** {pe['subject']}\n\n**תוכן:**\n{pe['body']}")
        if pe["attachments"]:
            st.caption("📎 קבצים מצורפים: " + ", ".join(Path(p).name for p in pe["attachments"]))

        col1, col2 = st.columns(2)
        if col1.button("✅ אשר ושלח", key="confirm_email"):
            result = _do_send_email(**pe)
            st.session_state.pending_email = None
            st.session_state.chat_history.append({"role": "assistant", "content": result})
            persist_current_conversation()
            st.rerun()
        if col2.button("❌ בטל", key="cancel_email"):
            st.session_state.pending_email = None
            st.session_state.chat_history.append({"role": "assistant", "content": "המייל בוטל ולא נשלח."})
            persist_current_conversation()
            st.rerun()


# ========================================================
# Chat input
# ========================================================

user_prompt = st.chat_input("כתוב הודעה...")

if user_prompt:
    full_prompt = user_prompt
    current_attachments = list(st.session_state.uploaded_paths)
    if current_attachments:
        names = ", ".join(Path(p).name for p in current_attachments)
        full_prompt += (
            f"\n\n(המשתמש צירף את הקבצים הבאים ואתה יכול לראות/לנתח את התוכן שלהם ישירות: {names}. "
            f"הנתיבים בדיסק לעיבוד נוסף: {', '.join(current_attachments)})"
        )

    st.session_state.chat_history.append({"role": "user", "content": full_prompt})

    with st.chat_message("user"):
        st.markdown(user_prompt)
        if st.session_state.uploaded_paths:
            st.caption("📎 " + ", ".join(Path(p).name for p in st.session_state.uploaded_paths))

    with st.chat_message("assistant"):
        response_placeholder = st.empty()
        response_placeholder.markdown("חושב...")

        try:
            assistant_reply = call_agent(st.session_state.chat_history, attachments=current_attachments)
            response_placeholder.markdown(assistant_reply)
        except Exception as e:
            assistant_reply = f"שגיאה: {str(e)}"
            response_placeholder.error(assistant_reply)

    st.session_state.chat_history.append({"role": "assistant", "content": assistant_reply})

    persist_current_conversation()

    st.rerun()
